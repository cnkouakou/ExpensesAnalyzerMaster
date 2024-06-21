# import required classes/functions/method 
from pptx import Presentation  
from pptx.chart.data import CategoryChartData  
from pptx.enum.chart import XL_CHART_TYPE  
from pptx.util import Inches 
  
  
# Create presentation object 
ppt = Presentation()  
  
# Adding slide with specific layout 
slide = ppt.slides.add_slide(ppt.slide_layouts[6]) 
  
# Define chart data  
# Creating object of chart 
chart_data = CategoryChartData()  
  
# Adding categories to chart 
chart_data.categories = ['East', 'West', 'Midwest']   
  
# Adding series 
chart_data.add_series('Series 1',  
                      (int(input("Enter Value:")),  
                        int(input("Enter Value:")), 
                        int(input("Enter Value:"))))  
  
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)  
  
slide.shapes.add_chart( XL_CHART_TYPE.COLUMN_CLUSTERED, x, 
                       y, cx, cy, chart_data ) 

# Adding a blank slide in out ppt 
slide = ppt.slides.add_slide(ppt.slide_layouts[6]) 
  
# Adjusting the width !   
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)  
  
# Adding tables 
shape = slide.shapes.add_table(3, 4, x,  
                               y, cx, cy)  
# Saving file 
ppt.save('chart-Tutorial.pptx') 
  
print("done")
