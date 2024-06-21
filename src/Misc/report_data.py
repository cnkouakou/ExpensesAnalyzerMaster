from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
import argparse
import pandas as pd
import numpy as np
from datetime import date
import matplotlib.pyplot as plt
import seaborn as sns

# Functions go here

def create_pivot(df, index_list=["Manager", "Rep", "Product"],
                 value_list=["Price", "Quantity"]):
    """
    Take a DataFrame and create a pivot table
    Return it as a DataFrame pivot table
    """
    table = pd.pivot_table(df, index=index_list,
                           values=value_list,
                           aggfunc=[np.sum, np.mean], fill_value=0)
    return table

def create_chart(df, filename):
    """ Create a simple bar chart saved to the filename based on the dataframe
    passed to the function
    """
    df['total'] = df['Quantity'] * df['Price']
    final_plot = df.groupby('Name')['total'].sum().order().plot(kind='barh')
    fig = final_plot.get_figure()
    # Size is the same as the PowerPoint placeholder
    fig.set_size_inches(6, 4.5)
    fig.savefig(filename, bbox_inches='tight', dpi=600)

def create_ppt(input, output, report_data, chart):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Quarterly Report"
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())

    # Create the summary graph
    graph_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Sales by account"
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(chart)
    subtitle = slide.placeholders[2]
    subtitle.text = "Results consistent with last quarter"


# Create a slide for each manager
for manager in report_data.index.get_level_values(0).unique():
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "Report for {}".format(manager)
    top = Inches(1.5)
    left = Inches(0.25)
    width = Inches(9.25)
    height = Inches(5.0)
    # Flatten the pivot table by resetting the index
    # Create a table on the slide
    df_to_table(slide, report_data.xs(manager, level=0).reset_index(),
                left, top, width, height)
prs.save(output)


if __name__ == "__main__":
    args = parse_args()
    df = pd.read_excel(args.report.name)
    report_data = create_pivot(df)
    create_chart(df, "report-image.png")
    create_ppt(args.infile.name, args.outfile.name, report_data, "report-image.png")