#https://python-pptx.readthedocs.io/en/latest/#community-guide
# import Presentation class 
# from pptx library 
from pptx import Presentation	
import pptx
from pptx.util import Inches 

import datetime
import os
from datetime import date
from datetime import datetime
from Helper import ___PATH___, ___FILE___LINE, log, level0,level1,level2,level3

""" Ref for slide types: 
0 -> title and subtitle 
1 -> title and content 
2 -> section header 
3 -> two content 
4 -> Comparison 
5 -> Title only 
6 -> Blank 
7 -> Content with caption 
8 -> Pic with caption 
"""
class PPTXReport():
    reportfilename = ''
    pptxcurrentmonth = ''
    pptxcurrentyear = 2024
    pptxtitle = ''
    pptxConfigfile = ''
    srcPath = ''
   
    def __init__(self, currentmonth, currentyear ) -> None:
        self.pptxcurrentmonth = currentmonth
        self.pptxcurrentyear = currentyear
        self.srcPath = ___PATH___()
        self.srcPath = self.srcPath.replace('src', '')
        self.pptxConfigfile = self.srcPath + 'configurations\\config.json'
        self.pptxReportingPath = self.srcPath + 'Reports\\'
        self.pptxTemplatePath = self.srcPath + 'Templates\\'
        self.pptxtitle = f"Household Financial Report: {self.pptxcurrentmonth} {self.pptxcurrentyear}"
    
    def getReportFilename(self):
        global reportfilename
        dtt = f'{datetime.now()}'
        for ele in [':', '-', '.', ' ']:
            dtt = dtt.replace(ele, '')
            self.reportfilename = self.pptxReportingPath + dtt + 'AnnualReport.pptx'


    def add_an_image_page(self, prs, img_path, strtitle):
        if os.path.isfile(img_path):
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            #slide.shapes.title.caption = strtitle
            left = top = Inches(1)
            pic = slide.shapes.add_picture(img_path, pptx.util.Inches(0.5), pptx.util.Inches(1.75),
                                width=pptx.util.Inches(9), height= pptx.util.Inches(5)) 
    
    def add_an_table_page(elf, prs, title, data):
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        title = slide.shapes.title
        title.text = title
        left = top = Inches(1)


    def generatepptxReport(self):
        self.pptxtitle = f"Household Financial Report: {self.pptxcurrentmonth} {self.pptxcurrentyear}"
        root = Presentation() 
        imagepath = f'{self.pptxReportingPath}ExpensePiePlotpng.png'
        self.add_an_image_page(root, imagepath, "First image page" )
        self.getReportFilename()
        root.save(self.reportfilename) 
        return self.reportfilename

    def generateAllpptxReport(self):
        self.pptxtitle = f"Household Financial Report: {self.pptxcurrentmonth} {self.pptxcurrentyear}"
        root = Presentation(f'{self.pptxTemplatePath}My_Template.pptx') 
        for file in os.listdir(self.pptxReportingPath):   
            # check the extension of files
            if file.endswith('.png'):
                imagepath = f'{self.pptxReportingPath}{file}'
                self.add_an_image_page(root, imagepath, "First image page" )
        self.getReportFilename()
        root.save(self.reportfilename) 
        return self.reportfilename